from fastapi import FastAPI, UploadFile, HTTPException, Form
from fastapi.staticfiles import StaticFiles
import fitz  # PyMuPDF
import docx  # python-docx
from pptx import Presentation  # python-pptx
import pandas as pd  # pandas
import gc
import io
import os
import base64
import re
import hashlib
from collections import Counter

# ---------------------------------------------------------------
# ตั้งค่า
# ---------------------------------------------------------------
IMAGE_DIR = "static/extracted_images"
os.makedirs(IMAGE_DIR, exist_ok=True)

# หมายเหตุ 5.0: ไม่เขียนไฟล์รูปลงดิสก์อีกแล้ว รูปถูกส่งกลับเป็น base64
# ให้ n8n เขียนลง C:\\n8n\\images\\ เอง ดู save_page_images()

# ข้ามรูปที่เล็กกว่านี้ (โลโก้ ไอคอน เส้นคั่น) หน่วยเป็น byte
MIN_IMAGE_SIZE = 15000

# หัวข้อ = ฟอนต์ใหญ่กว่าเนื้อความปกติกี่เท่า
HEADING_SIZE_RATIO = 1.15

# ความยาวสูงสุดต่อ 1 section ถ้าเกินจะถูกซอยย่อย
MAX_SECTION_CHARS = 1200

# ถ้า section สั้นกว่านี้ จะถูกรวมกับ section ถัดไป
# หมายเหตุ: อย่าตั้งสูงเกินไป เพราะเคสที่เนื้อหาสั้นแต่จบในตัวเอง
# จะถูกลากไปรวมกับเคสถัดไป กลายเป็น chunk ที่ปนสองเรื่อง
MIN_SECTION_CHARS = 80

# ถ้า section มีคำเหล่านี้อยู่แล้ว ถือว่า "จบในตัวเอง" ห้ามลากไปรวมกับอันถัดไป
# แม้เนื้อหาจะสั้นกว่า MIN_SECTION_CHARS ก็ตาม
# (ปรับตามรูปแบบคู่มือของแต่ละองค์กร ถ้าไม่ต้องการให้ตั้งเป็น [] )
SELF_CONTAINED_MARKERS = ["Resolution Steps"]

app = FastAPI(
    title="Multi-Format Document Extraction API",
    description="สกัดข้อความและรูปจากเอกสาร แบ่ง chunk ตามหัวข้อ สำหรับ Vector DB / RAG",
    version="5.2.0",
)

app.mount("/static", StaticFiles(directory="static"), name="static")


@app.get("/")
def root():
    return {"message": "Document Extraction Service is running perfectly!"}


@app.get("/health")
def health_check():
    return {"status": "running"}


# ===============================================================
# ส่วนที่ 0 : ชื่อระบบจากชื่อไฟล์ และภาษาของเนื้อหา
# ===============================================================

def guess_system_name(filename):
    """
    เดาชื่อระบบจากชื่อไฟล์

    [แก้บั๊ก] ของเดิม: filename.split("_")[0] if "_" in filename else "unknown"
    ทำให้ไฟล์ที่ไม่มีขีดล่างกลายเป็น unknown ทั้งหมด
      ACCA_manual_pdf.pdf  ->  ACCA        (บังเอิญถูก)
      EPayment.pdf         ->  unknown     (ผิด)

    ของใหม่: ตัดนามสกุลออกก่อน แล้วเอาส่วนหน้าขีดล่างถ้ามี
             ถ้าไม่มีขีดล่างก็ใช้ชื่อไฟล์ทั้งอัน ไม่ตกไปเป็น unknown อีก
      ACCA_manual_pdf.pdf  ->  ACCA
      EPayment.pdf         ->  EPayment
    """
    stem = os.path.splitext(filename)[0].strip()
    if not stem:
        return "unknown"
    return stem.split("_")[0] if "_" in stem else stem


def detect_lang(text):
    """
    [เพิ่มใน 5.2] ตรวจภาษาของ chunk เพื่อติดป้าย lang ลง vector metadata

    ใช้แยก tool ค้นหาใน n8n เป็นสองตัว (lang=th / lang=en) โดยให้แต่ละตัว
    มีโควตาผลลัพธ์ของตัวเอง ป้องกันไม่ให้ chunk ภาษาเดียวกับคำถาม
    กวาดที่นั่งไปหมดจน chunk อีกภาษาที่ตอบตรงกว่าไม่ถูกดึงมาเลย

    กฎ: เจออักขระไทยแม้ตัวเดียว = th, เป็นอังกฤษล้วน = en

    ⚠️ สำคัญ: ต้องเรียกด้วย "เนื้อหาล้วน" เท่านั้น ห้ามส่งข้อความที่เติม
    header แล้วเข้ามา เพราะ header ทุกอันมีคำว่า "หน้า" และ "หัวข้อ"
    ซึ่งเป็นภาษาไทย จะทำให้ทุก chunk ถูกตัดสินเป็น th หมดทั้งฐาน
    """
    return "th" if re.search(r"[ก-๙]", text or "") else "en"


# ===============================================================
# ส่วนที่ 1 : หาว่าอะไรคือ "หัวข้อ"
# ===============================================================

def get_body_font_size(doc):
    """
    หาขนาดฟอนต์ของเนื้อความปกติ โดยดูว่าขนาดไหนถูกใช้มากที่สุดทั้งเล่ม
    ใช้เป็นฐานเทียบว่าอะไรคือหัวข้อ
    """
    counter = Counter()

    for page in doc:
        data = page.get_text("dict")
        for block in data.get("blocks", []):
            if block.get("type") != 0:      # 0 = text block
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if not text:
                        continue
                    size = round(span.get("size", 0), 1)
                    counter[size] += len(text)   # ถ่วงน้ำหนักด้วยจำนวนตัวอักษร

    if not counter:
        return 11.0
    return counter.most_common(1)[0][0]


def is_heading(line, body_size):
    """
    ตัดสินว่าบรรทัดนี้เป็นหัวข้อหรือไม่
    เกณฑ์: ฟอนต์ใหญ่กว่าปกติ หรือ ตัวหนา + สั้น

    [แก้บั๊ก 4.2] เอากฎ "ขึ้นต้นด้วยเลขข้อ" ออก
    กฎเดิมทำให้ขั้นตอน 1. 2. 3. 4. ในเนื้อหา ถูกมองเป็นหัวข้อใหม่ทุกข้อ
    ผลคือขั้นตอนแก้ปัญหาชุดเดียวถูกหั่นเป็น chunk ละข้อ
    คำถามอยู่ chunk หนึ่ง คำตอบอยู่อีก chunk หนึ่ง

    [แก้บั๊ก 4.3] กรองเลขลอย ๆ และเลขโรมันออก
    พบใน EPayment.pdf ว่า 153 จาก 195 chunk (78%) มี heading เป็น
    ตัวเลขล้วนหรือเลขโรมันล้วน เช่น "3" "2" "II." "III."
    เพราะเลขข้อหรือเลขหน้าที่พิมพ์ตัวหนา เข้าเงื่อนไข is_bold + สั้น
    ทำให้ chunk ถูกหั่นทุกครั้งที่เจอเลข และ heading ไม่ให้บริบทกับ LLM เลย
    """
    spans = line.get("spans", [])
    if not spans:
        return False

    text = "".join(s.get("text", "") for s in spans).strip()
    if not text or len(text) > 120:      # หัวข้อมักไม่ยาวมาก
        return False

    # --- กรองสิ่งที่ไม่มีทางเป็นหัวข้อ ---

    # ต้องมีตัวอักษรอย่างน้อย 1 ตัว (ไทยหรืออังกฤษ)
    # กรอง "3" "2." "1)" "-" "•" ที่เป็นเลขข้อหรือ bullet
    if not re.search(r"[A-Za-zก-๙]", text):
        return False

    # เลขโรมันล้วน เช่น "II." "III." "IV)" "i."
    # มีตัวอักษรจึงรอดกฎข้างบนมาได้ แต่ไม่ใช่หัวข้อ
    if re.fullmatch(r"[IVXLCDMivxlcdm]+[\.\)]?", text):
        return False

    # สั้นเกินกว่าจะเป็นหัวข้อที่มีความหมาย
    if len(text) < 3:
        return False

    # --- เกณฑ์ตัดสินว่าเป็นหัวข้อ ---

    max_size = max(s.get("size", 0) for s in spans)
    # flags bit 4 (ค่า 16) = ตัวหนา
    is_bold = any(s.get("flags", 0) & 16 for s in spans)

    if max_size >= body_size * HEADING_SIZE_RATIO:
        return True

    # ตัวหนา + สั้น = น่าจะเป็นหัวข้อ
    if is_bold and len(text) < 80:
        return True

    return False


def is_toc_line(text):
    """
    บรรทัดสารบัญ เช่น  Trip Activity ..................... 5
    ตรวจจากจุดไข่ปลาที่ต่อกันยาว ๆ
    """
    return bool(re.search(r"\.{6,}", text))


# ===============================================================
# ส่วนที่ 2 : เซฟรูป พร้อมจำตำแหน่ง y ไว้
# ===============================================================

def save_page_images(doc, page, page_num):
    """
    สกัดรูปในหน้านั้น คืน list ของ dict ที่มี filename, base64 และตำแหน่ง y
    ตำแหน่ง y ใช้จับคู่รูปกับหัวข้อที่อยู่ใกล้กัน

    [เปลี่ยนใน 5.0] ไม่เขียนไฟล์ลงดิสก์อีกแล้ว
    เดิมเขียนลง static/extracted_images บน Render ซึ่งเป็น ephemeral filesystem
    ไฟล์หายทุกครั้งที่ sleep / deploy / ย้าย container
    ทำให้ลิงก์ "ดูภาพประกอบ" ที่ผู้ใช้กดขึ้น 404 (ปัญหาข้อ 2.1)

    ของใหม่: คืนรูปเป็น base64 ให้ n8n เอาไปเขียนลง C:\\n8n\\images\\ เอง
    ซึ่งเป็นดิสก์จริงบนเซิร์ฟเวอร์บริษัท ไฟล์ไม่หาย

    [แก้บั๊ก 5.1] เปลี่ยนชื่อไฟล์จาก uuid สุ่ม เป็น hash ของเนื้อรูป
    ของเดิม: f"page_{page_num}_{uuid.uuid4().hex[:8]}.{ext}"
    uuid4() สุ่มใหม่ทุกครั้ง ทำให้ ingest ซ้ำได้ชื่อไฟล์ใหม่หมด
    ผลคือชื่อที่เก็บใน vector metadata (รอบเก่า) ไม่ตรงกับไฟล์บนดิสก์ (รอบใหม่)
    n8n อ่านไฟล์ไม่เจอ ขึ้น "No file(s) found"

    ของใหม่: ใช้ sha256 ของ image_bytes รูปเดิมได้ชื่อเดิมเสมอ
    ingest ซ้ำกี่รอบก็ตรงกัน และรูปซ้ำข้ามหน้าก็ยังแยกกันด้วย page_num
    """
    saved = []

    for img in page.get_images(full=True):
        xref = img[0]

        try:
            base_image = doc.extract_image(xref)
        except Exception:
            continue

        image_bytes = base_image.get("image")
        image_ext = base_image.get("ext", "png")

        if not image_bytes or len(image_bytes) < MIN_IMAGE_SIZE:
            continue

        # หาว่ารูปนี้วางอยู่ตรงไหนในหน้า
        try:
            rects = page.get_image_rects(xref)
            y_top = rects[0].y0 if rects else 0.0
            y_bottom = rects[0].y1 if rects else 0.0
        except Exception:
            y_top, y_bottom = 0.0, 0.0

        # ชื่อไฟล์ต้องคงที่ ห้ามสุ่ม ไม่งั้น metadata กับดิสก์จะไม่ตรงกัน
        img_hash = hashlib.sha256(image_bytes).hexdigest()[:8]
        filename = f"page_{page_num}_{img_hash}.{image_ext}"

        saved.append({
            "filename": filename,
            "base64": base64.b64encode(image_bytes).decode("ascii"),
            "mime": f"image/{'jpeg' if image_ext in ('jpg', 'jpeg') else image_ext}",
            "size": len(image_bytes),
            "y_top": y_top,
            "y_bottom": y_bottom,
        })

    return saved


# ===============================================================
# ส่วนที่ 3 : แบ่งเอกสารเป็น section ตามหัวข้อ
# ===============================================================

def extract_sections(doc):
    """
    เดินทีละหน้า ทีละบรรทัด แล้วรวมเป็น section
    section ใหม่เริ่มเมื่อเจอหัวข้อ และไหลข้ามหน้าได้

    หมายเหตุ: page_lines.sort(key=lambda x: x["y"]) ถูกต้องแล้ว
    PyMuPDF นับ y จากบนลงล่าง เรียงน้อยไปมาก = บนลงล่าง ห้ามกลับทิศ
    """
    body_size = get_body_font_size(doc)
    sections = []
    all_images = []          # เก็บ base64 ของรูปทุกใบ ส่งกลับให้ n8n เขียนลงดิสก์

    current = {
        "heading": None,
        "lines": [],
        "page_start": 1,
        "page_end": 1,
        "images": [],
    }

    for page_index, page in enumerate(doc):
        page_num = page_index + 1
        page_images = save_page_images(doc, page, page_num)
        used_images = set()

        data = page.get_text("dict")

        # รวบรวมทุกบรรทัดในหน้า พร้อมตำแหน่ง y แล้วเรียงจากบนลงล่าง
        page_lines = []
        for block in data.get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                text = "".join(s.get("text", "") for s in line.get("spans", [])).strip()
                if not text:
                    continue
                # ข้ามบรรทัดสารบัญ ไม่มีประโยชน์ต่อการค้นหา
                if is_toc_line(text):
                    continue
                y = line.get("bbox", [0, 0, 0, 0])[1]
                page_lines.append({"text": text, "y": y, "raw": line})

        page_lines.sort(key=lambda x: x["y"])

        for idx, item in enumerate(page_lines):
            if is_heading(item["raw"], body_size):
                # ปิด section เดิมก่อนเริ่มอันใหม่
                if current["lines"] or current["heading"]:
                    sections.append(current)

                current = {
                    "heading": item["text"],
                    "lines": [],
                    "page_start": page_num,
                    "page_end": page_num,
                    "images": [],
                }
            else:
                current["lines"].append(item["text"])
                current["page_end"] = page_num

            # ผูกรูปที่อยู่ระหว่างบรรทัดนี้กับบรรทัดถัดไป เข้ากับ section ปัจจุบัน
            y_now = item["y"]
            y_next = page_lines[idx + 1]["y"] if idx + 1 < len(page_lines) else 10_000

            for img_i, img in enumerate(page_images):
                if img_i in used_images:
                    continue
                if y_now <= img["y_top"] < y_next:
                    current["images"].append(img["filename"])
                    used_images.add(img_i)

        # รูปที่ยังจับคู่ไม่ได้ (เช่นอยู่บนสุดของหน้า) ให้ผูกกับ section ปัจจุบัน
        for img_i, img in enumerate(page_images):
            if img_i not in used_images:
                current["images"].append(img["filename"])

        # เก็บข้อมูลรูปทั้งหมดไว้ที่เดียว เพื่อส่งกลับให้ n8n เขียนลงดิสก์
        for img in page_images:
            all_images.append({
                "filename": img["filename"],
                "base64": img["base64"],
                "mime": img["mime"],
                "size": img["size"],
                "page": page_num,
            })

    if current["lines"] or current["heading"]:
        sections.append(current)

    return sections, all_images


def _is_self_contained(section):
    """
    section นี้จบในตัวเองหรือยัง ถ้าใช่จะไม่ถูกลากไปรวมกับอันถัดไป
    กันเคสที่เนื้อหาสั้นแต่ครบชุด ถูกรวมกับเรื่องอื่นจนปนกัน
    """
    if not SELF_CONTAINED_MARKERS:
        return False
    body = "\n".join(section["lines"])
    return any(marker in body for marker in SELF_CONTAINED_MARKERS)


def merge_short_sections(sections):
    """
    section ที่สั้นเกินไป (เช่นหัวข้อลอย ๆ ไม่มีเนื้อหา) ให้รวมกับอันถัดไป

    [แก้บั๊กใหญ่] ของเดิมเขียนว่า
        sec["lines"] = [head] + buffer["lines"] + sec["lines"]
    คือเอา section ปัจจุบันไปแปะ "หน้า" section ถัดไป
    พอ section สั้นต่อกันหลายอัน เนื้อหาจะเรียงถอยหลังสะสม
    และ heading ที่เหลือกลายเป็นอันล่างสุดของหน้า

    ตัวอย่างหน้าปก ACCA ที่เคยได้ออกมา (กลับหัวทั้งหมด)
        heading: Vaneenut Singkanthana
        lines:   Create By / Version Common Issues /
                 Advance Claim Application (ACCA) / User Manual Document

    ของใหม่: ต่อท้ายตามลำดับจริง และเก็บ heading "แรก" ไว้เป็นชื่อ section
    """
    merged = []
    buffer = None

    for sec in sections:
        if buffer is None:
            buffer = sec
            continue

        buffer_body = "\n".join(buffer["lines"]).strip()
        too_short = len(buffer_body) < MIN_SECTION_CHARS

        if too_short and not _is_self_contained(buffer):
            lines = list(buffer["lines"])

            # หัวข้อของ section ถัดไป กลายเป็นบรรทัดเนื้อหาต่อท้าย
            # (ไม่ใช่แปะไว้ข้างหน้าเหมือนโค้ดเดิม)
            if sec["heading"]:
                lines.append(sec["heading"])

            lines.extend(sec["lines"])

            buffer = {
                # เก็บหัวข้อแรกไว้ ถ้าอันแรกไม่มีค่อยใช้ของอันถัดไป
                "heading": buffer["heading"] or sec["heading"],
                "lines": lines,
                "page_start": buffer["page_start"],
                "page_end": sec["page_end"],
                "images": buffer["images"] + sec["images"],
            }
        else:
            merged.append(buffer)
            buffer = sec

    if buffer is not None:
        merged.append(buffer)

    return merged


def split_long_section(section):
    """
    section ที่ยาวเกินไป ให้ซอยเป็นชิ้นย่อย
    แต่ยังคงใส่หัวข้อเดิมนำหน้าทุกชิ้น เพื่อไม่ให้หลุดบริบท
    (การใส่หัวข้อนำหน้าทำใน endpoint ตอนประกอบ header)
    """
    body = "\n".join(section["lines"]).strip()

    if len(body) <= MAX_SECTION_CHARS:
        return [body]

    parts = []
    buf = ""

    for line in section["lines"]:
        if len(buf) + len(line) + 1 > MAX_SECTION_CHARS and buf:
            parts.append(buf.strip())
            buf = line
        else:
            buf = f"{buf}\n{line}" if buf else line

    if buf.strip():
        parts.append(buf.strip())

    return parts


# ===============================================================
# ส่วนที่ 4 : endpoint
# ===============================================================

@app.post("/extract")
async def extract_file(
    file: UploadFile,
    system: str = Form(""),
    doc_type: str = Form(""),
):
    try:
        content = await file.read()
        file_hash = hashlib.sha256(content).hexdigest()
        results = []
        extracted_images = []
        filename = file.filename
        lower_name = filename.lower()

        # ถ้าไม่ได้ส่งมา ให้เดาจากชื่อไฟล์
        if not system:
            system = guess_system_name(filename)
        if not doc_type:
            doc_type = "faq" if "faq" in lower_name else "manual"

        # -------------------- PDF --------------------
        if lower_name.endswith(".pdf"):
            doc = fitz.open(stream=content, filetype="pdf")

            sections, extracted_images = extract_sections(doc)
            sections = merge_short_sections(sections)

            for sec in sections:
                heading = sec["heading"] or "(ไม่มีหัวข้อ)"
                page_start = sec["page_start"]
                page_end = sec["page_end"]
                page_label = (
                    f"หน้า {page_start}"
                    if page_start == page_end
                    else f"หน้า {page_start}-{page_end}"
                )

                for part_index, part in enumerate(split_long_section(sec)):
                    if not part.strip():
                        continue

                    # ใส่ path รูปไว้ในตัวข้อความด้วย
                    # เพราะ AI Agent มองไม่เห็น metadata เห็นแค่ text
                    # ใส่เฉพาะชิ้นแรก กันรูปซ้ำหลาย chunk
                    img_line = ""
                    if sec["images"] and part_index == 0:
                        img_line = "[รูปประกอบ: " + ", ".join(sec["images"]) + "]\n"

                    header = (
                        f"[{system} - {doc_type} - {page_label}]\n"
                        f"[หัวข้อ: {heading}]\n"
                        f"{img_line}"
                    )

                    results.append({
                        "page": page_start,
                        "page_end": page_end,
                        "heading": heading,
                        "part": part_index + 1,
                        "text": header + part,
                        # ตรวจจาก heading + เนื้อหาเท่านั้น ไม่รวม header
                        # เพราะ header มีคำว่า หน้า/หัวข้อ เป็นไทยทุกอัน
                        "lang": detect_lang(f"{heading}\n{part}"),
                        "images": sec["images"] if part_index == 0 else [],
                    })

            doc.close()

        # -------------------- ไฟล์ข้อความ --------------------
        elif lower_name.endswith((".txt", ".md", ".json", ".log", ".csv")):
            raw_text = content.decode("utf-8", errors="ignore")
            results.append({
                "page": 1,
                "heading": filename,
                "text": f"[{system} - {doc_type}]\n{raw_text}",
                "lang": detect_lang(raw_text),
                "images": [],
            })

        # -------------------- Word --------------------
        elif lower_name.endswith(".docx"):
            d = docx.Document(io.BytesIO(content))

            current_heading = "(ไม่มีหัวข้อ)"
            buf = []

            def flush():
                if buf:
                    body = "\n".join(buf)
                    results.append({
                        "page": 1,
                        "heading": current_heading,
                        "text": f"[{system} - {doc_type}]\n[หัวข้อ: {current_heading}]\n"
                                + body,
                        "lang": detect_lang(f"{current_heading}\n{body}"),
                        "images": [],
                    })

            for p in d.paragraphs:
                text = p.text.strip()
                if not text:
                    continue

                # ใช้ style ของ Word บอกว่าเป็นหัวข้อ
                if p.style.name.startswith("Heading") or p.style.name == "Title":
                    flush()
                    buf = []
                    current_heading = text
                else:
                    buf.append(text)

            flush()

        # -------------------- PowerPoint --------------------
        elif lower_name.endswith(".pptx"):
            prs = Presentation(io.BytesIO(content))
            for idx, slide in enumerate(prs.slides):
                slide_text = [
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
                title = slide_text[0] if slide_text else f"สไลด์ {idx + 1}"
                body = "\n".join(slide_text)
                results.append({
                    "page": idx + 1,
                    "heading": title,
                    "text": f"[{system} - {doc_type} - สไลด์ {idx + 1}]\n[หัวข้อ: {title}]\n"
                            + body,
                    "lang": detect_lang(body),
                    "images": [],
                })

        # -------------------- Excel --------------------
        elif lower_name.endswith((".xlsx", ".xls")):
            excel_file = pd.ExcelFile(io.BytesIO(content))
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                csv_text = df.to_csv(index=False)
                results.append({
                    "page": sheet_name,
                    "heading": sheet_name,
                    "text": f"[{system} - {doc_type} - Sheet: {sheet_name}]\n"
                            + csv_text,
                    # ชื่อชีตมักเป็นไทย แต่เนื้อตารางอาจเป็นอังกฤษล้วน
                    # ตรวจจากเนื้อตารางเป็นหลัก เพราะนั่นคือสิ่งที่ผู้ใช้ค้นหา
                    "lang": detect_lang(csv_text),
                    "images": [],
                })

        else:
            raise HTTPException(
                status_code=400,
                detail=f"ไม่รองรับไฟล์ประเภทนี้: {filename}",
            )

        del content
        gc.collect()

        # นับจำนวน chunk แต่ละภาษา ไว้เช็คว่าแยกภาษาถูกต้องไหม
        # ถ้าเห็น th 100% ให้สงสัยว่ามีการส่ง header เข้า detect_lang()
        lang_summary = {"th": 0, "en": 0}
        for r in results:
            lang_summary[r.get("lang", "en")] = lang_summary.get(r.get("lang", "en"), 0) + 1

        return {
            "status": "ok",
            "filename": filename,
            "file_hash": file_hash,
            "system": system,
            "doc_type": doc_type,
            "total_pages": len(results),
            "lang_summary": lang_summary,
            "pages": results,
            # รูปทั้งหมดเป็น base64 ให้ n8n เขียนลง C:\n8n\images\ เอง
            # ชื่อไฟล์ตรงกับที่ปรากฏใน [รูปประกอบ: ...] ของแต่ละ chunk
            "images": extracted_images,
            "image_count": len(extracted_images),
        }

    except HTTPException as he:
        raise he
    except Exception as e:
        gc.collect()
        raise HTTPException(
            status_code=500,
            detail=f"เกิดข้อผิดพลาดในการประมวลผลไฟล์: {str(e)}",
        )